import os
from io import BytesIO

import pandas as pd
import requests
import streamlit as st
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

load_dotenv()
FMP_API_KEY = os.getenv("FMP_API_KEY")
STOCKS = [
    # "AAPL",
    # "AL",
    # "AX",
    # "DAR",
    # "INMD",
    "META",
    "MSFT",
    "NVDA",
    # "PAYC",
    # "SCHW",
    # "SWKS",
    # "TSCO",
    # "V",
    # "VRTX",
]


### Functions ###
@st.cache_data
def get_last_close_price(ticker: str):
    """Get last closing price from Polygon.io"""

    res = requests.get(
        f"https://financialmodelingprep.com/api/v3/quote/{ticker}",
        params={"apikey": FMP_API_KEY},
    )
    if res.status_code != 200:
        st.warning(f"Failed to fetch data from Polygon ({res.status_code}): {res.text}")
        return
    data = res.json()
    return data[0]["previousClose"]


@st.cache_data
def get_financial_reports_fmp(ticker: str):
    """Get financial reports from Polygon.io"""
    res = requests.get(
        f"https://financialmodelingprep.com/api/v3/income-statement/{ticker}",
        params={"apikey": FMP_API_KEY, "period": "quarter"},
    )
    if res.status_code != 200:
        st.warning(f"Failed to fetch data from FMP ({res.status_code}): {res.text}")
        return
    data = res.json()
    return data


def get_financials_df(ticker: str) -> pd.DataFrame:
    quarterly_financials = get_financial_reports_fmp(ticker)
    df = pd.DataFrame(quarterly_financials)
    df = df.rename(
        columns={
            "calendarYear": "year",
            "period": "quarter",
            "incomeBeforeTax": "income",
            "epsdiluted": "eps",
        }
    )
    df = df[["year", "quarter", "revenue", "income", "eps"]]
    df = df.loc[:,~df.columns.duplicated()]
    df = df.sort_values(["year", "quarter"], ascending=False)
    df = df.set_index(["year", "quarter"])
    yoy_change = (df - df.shift(-4)) / df.shift(-4) * 100
    yoy_change = yoy_change.dropna().rename(columns=lambda x: f"{x}_yoy_change")
    df = pd.concat([df, yoy_change], axis=1)
    df["period"] = df.index.get_level_values("quarter") + df.index.get_level_values(
        "year"
    ).astype(str)
    df["pre_tax_profit_margin"] = df["income"] / df["revenue"] * 100
    return df


def make_bar_plot(
    df: pd.DataFrame, col: str, title: str, color="green", target=7.5, n_quarters=4, ax=None
):
    ax = df.iloc[:n_quarters][::-1].plot.bar(x="period", y=col, color=color, rot=0, ax=ax)
    ax.grid(axis="y", color="gray", linestyle="-", linewidth=0.5, alpha=0.2)
    ax.get_legend().remove()
    # Remove lines
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.set_xlabel("")
    ax.set_title(title)
    ax.axhline(target, color="red", linewidth=2.0, linestyle="--")
    ax.tick_params(length=0)
    ax.set_ylim()
    ax.set_yticklabels([f"{int(x)}%" for x in ax.get_yticks()])
    for container in ax.containers:
        ax.bar_label(container, fmt="%.1f%%", label_type="edge")
    return ax


def make_ranges_plot(
    current_price: float,
    buy_lower_price: float,
    hold_lower_price: float,
    sell_lower_price: float,
    sell_upper_price: float,
    ax=None
):
    ranges = pd.DataFrame(
        [
            [buy_lower_price, hold_lower_price - buy_lower_price],
            [hold_lower_price, sell_lower_price - hold_lower_price],
            [sell_lower_price, sell_upper_price - sell_lower_price],
        ],
        columns=["Low", "High"],
        index=["Buy", "Hold", "Sell"],
    )
    ax = ranges.plot.bar(stacked=True, color=[(0, 0, 0, 0), "grey"], rot=0, ax=ax)
    ax.get_legend().remove()
    ax.grid(axis="y", color="gray", linestyle="-", linewidth=0.5, alpha=0.2)
    for container in ax.containers:
        ax.bar_label(container, fmt="$%.1f", label_type="edge")
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(length=0)
    ax.axhline(current_price, color="black", linewidth=0.5)
    ax.text(
        2,
        current_price,
        f"Last close: (${current_price:.2f})",
        va="center",
        ha="center",
        backgroundcolor=(1, 1, 1, 0.7),
    )
    return ax


def make_ppt_report(
    ticker: str,
    author: str,
    financial_period: str,
    decision: str,
    comments: str,
    ax_revenue,
    ax_eps,
    ax_ptpm,
    ax_ranges,
):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"{ticker} {financial_period} Financial Report"
    title.text_frame.paragraphs[0].font.size = Inches(0.3)
    # Left justify title
    title.text_frame.paragraphs[0].alignment = 1
    # Move title up
    title.top = Inches(0.15)
    title.width = Inches(6)
    title.left = Inches(0.1)

    # Add author
    if author:
        text_box = slide.shapes.add_textbox(
            Inches(0.1), Inches(0.3), Inches(2), Inches(0.2)
        ).text_frame
        text_box.text = f"Report created by: {author}"
        text_box.paragraphs[0].font.size = Inches(0.15)

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(3)

    # Add recommendation
    text_box = slide.shapes.add_textbox(
        Inches(0.1), Inches(0.6), width * 2, Inches(0.5)
    ).text_frame
    text_box.text = f"Recommendation: {decision}"
    text_box.paragraphs[0].font.size = Inches(0.15)

    # Add comments
    text_box = slide.shapes.add_textbox(
        Inches(0.1), Inches(1.0), width * 2, Inches(0.5)
    ).text_frame
    text_box.text = comments
    text_box.paragraphs[0].font.size = Inches(0.1)
    text_box.word_wrap = True

    img_stream = BytesIO()
    ax_revenue.figure.savefig(img_stream, format="png")
    slide.shapes.add_picture(img_stream, left, top, width, height)

    img_stream = BytesIO()
    ax_eps.figure.savefig(img_stream, format="png")
    slide.shapes.add_picture(img_stream, left + width, top, width, height)

    img_stream = BytesIO()
    ax_ptpm.figure.savefig(img_stream, format="png")
    slide.shapes.add_picture(img_stream, left, top + height, width, height)

    img_stream = BytesIO()
    ax_ranges.figure.savefig(img_stream, format="png")
    slide.shapes.add_picture(img_stream, left + width, top + height, width, height)
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    return ppt_buffer

def make_plots(
    df: pd.DataFrame,
    ticker: str,
    current_price: float,
    revenue_target: float,
    ptpm_target: float,
    eps_target: float,
    buy: list[float],
    hold: list[float],
    sell: list[float],
    make_multipanel: bool = True,
    author: str=None,
    decision: str=None,
    financial_period: str=None,
    comments: str=None,
):
    ax = None
    if make_multipanel:
        fig, axes = plt.subplots(3,2, figsize=(10,10), height_ratios=(1,5,5))
        title = f"{ticker} {financial_period} Financial Report"
        fig.suptitle(title)
        axes[0,0].axis("off")
        axes[0,1].axis("off")
        subtitle = f"Report created by {author} | Recommendation: {decision}"
        fig.text(0.06, 0.9, subtitle, fontsize=12)
        fig.text(0.06, 0.85, comments, fontsize=10)
        fig.subplots_adjust(wspace=0.3, hspace=0.3)

    if make_multipanel:
        ax = axes[1,0]
    ax_revenue = make_bar_plot(
        df,
        "revenue_yoy_change",
        "Revenue YoY Growth Rate (%)",
        color="green",
        target=revenue_target,
        ax=ax
    )
    if not make_multipanel:
        st.pyplot(ax_revenue.figure)

    # EPS YoY
    if make_multipanel:
        ax = axes[1,1]
    ax_eps = make_bar_plot(
        df,
        "eps_yoy_change",
        "EPS YoY Growth Rate (%)",
        color="blue",
        target=eps_target,
        ax=ax
    )
    if not make_multipanel:
        st.pyplot(ax_eps.figure)

    # Pre-tax profit margin
    if make_multipanel:
        ax = axes[2,0]
    ax_ptpm = make_bar_plot(
        df,
        "pre_tax_profit_margin",
        "Pre-tax Profit Margin (%)",
        color="olive",
        target=ptpm_target,
        ax=ax
    )
    if not make_multipanel:
        st.pyplot(ax_ptpm.figure)

    # Buy, hold sell
    if make_multipanel:
        ax = axes[2,1]
    ax_ranges = make_ranges_plot(
        current_price, buy[0], hold[0], sell[0], sell[1],ax=ax
    )
    ax_ranges.set_title("Buy, Hold, Sell Ranges")
    if not make_multipanel:
        st.pyplot(ax_ranges.figure)
    if make_multipanel:
        return fig


### App ###
st.set_page_config(
    page_title="Stock Tracking Report",
    page_icon=":chart_with_upwards_trend:",
    layout="wide",
)
st.title("Stock Tracking Report")
input_panel, report_panel = st.columns([2, 1])

with input_panel:
    ticker = st.selectbox("Ticker", STOCKS, index=None)
    if not ticker:
        st.stop()
    df = get_financials_df(ticker)
    current_price = get_last_close_price(ticker)

    # Revenue, EPS and PTPM targets
    st.write("### 1. Targets")
    cols = st.columns(3)
    with cols[0]:
        revenue_target = st.number_input("Revenue YoY % Target", 0, value=10)
    with cols[1]:
        ptpm_target = st.number_input("PTPM (%) Target", 0, value=10)
    with cols[2]:
        eps_target = st.number_input("EPS YoY (%) Target", 0, value=10)

    # Buy sell, hold range
    st.write("")
    st.write("### 2. Buy, Sell, Hold Range")

    st.write("**Buy Range**")
    buy_cols = st.columns(4)
    with buy_cols[0]:
        buy_lower = st.number_input("Lower", 0, value=100, key="buy_lower")
    with buy_cols[1]:
        buy_upper = st.number_input("Upper", 0, value=150, key="buy_upper")

    st.write("")
    st.write("**Hold Range**")
    hold_cols = st.columns(4)
    with hold_cols[0]:
        hold_lower = st.number_input(
            "Lower",
            0,
            value=buy_upper,
            key="hold_lower",
            disabled=True,
        )
    with hold_cols[1]:
        hold_upper = st.number_input("Upper", 0, value=200, key="hold_upper")

    st.write("")
    st.write("**Sell Range**")
    sell_cols = st.columns(4)
    with sell_cols[0]:
        sell_lower = st.number_input(
            "Lower",
            0,
            value=hold_upper,
            key="sell_lower",
            disabled=True,
        )
    with sell_cols[1]:
        sell_upper = st.number_input("Upper", 0, value=250, key="sell_upper")

    # Comment
    st.write("### 3. Recommendation")
    decision = st.selectbox("My recommendation", ["Buy", "Sell", "Hold"])

    comments = st.text_area(
        "Comments", placeholder=f"This is a {decision.lower()} because..."
    )
    author = st.text_input("What is your name?")

    download_container = st.container()

with report_panel.container(border=True):
    # Create report
    st.write("_Report preview_")
    st.write(f"#### {ticker} {df.iloc[0]['period']} Report")
    if author:
        st.write(f"Report created by: **{author}**")
    st.write(f"Recommendation: **{decision}**")
    st.caption(comments)
    if df is not None:
        make_plots(
            df,
            ticker=ticker,
            current_price=current_price,
            revenue_target=revenue_target,
            ptpm_target=ptpm_target,
            eps_target=eps_target,
            buy=[buy_lower, buy_upper],
            hold=[hold_lower, hold_upper],
            sell=[sell_lower,sell_upper],
            make_multipanel=False
        )
    else:
        st.write(f"No data found for {ticker}. ")
with download_container:
    # report_buffer = make_ppt_report(
    #     ticker=ticker,
    #     author=author,
    #     decision=decision,
    #     financial_period=df.iloc[0]["period"],
    #     comments=comments,
    #     ax_revenue=ax_revenue,
    #     ax_eps=ax_eps,
    #     ax_ptpm=ax_ptpm,
    #     ax_ranges=ax_ranges,
    # )
    financial_period = df.iloc[0]["period"]
    fig = make_plots(
        df,
        ticker=ticker,
        current_price=current_price,
        revenue_target=revenue_target,
        ptpm_target=ptpm_target,
        eps_target=eps_target,
        buy=[buy_lower, buy_upper],
        hold=[hold_lower, hold_upper],
        sell=[sell_lower,sell_upper],
        make_multipanel=True,
        author=author,
        decision=decision,
        financial_period=financial_period,
        comments=comments,
    )
    pdf_buffer = BytesIO()
    fig.savefig(pdf_buffer, format="pdf")

    st.download_button(
        "Download report", pdf_buffer, f"Stock Report {ticker} {financial_period}.pdf", ".pdf", type="primary"
    )
